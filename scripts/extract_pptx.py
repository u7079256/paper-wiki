"""PPTX -> Markdown (+ embedded images) extraction, mirroring mineru's layout.

LOCAL FALLBACK ONLY. The preferred path for slide decks is to convert PPTX to PDF
on the GPU server (`soffice --headless --convert-to pdf`) and run the normal
mineru pipeline -- that preserves layout and gives the same OCR quality as the
other slides. Use this script only when no LibreOffice is available.

python-pptx can read text frames, tables, embedded image blobs and (1.0+) convert
OOXML math to LaTeX -- but slide LAYOUT / positioning / image-to-caption relations
are LOST, and SmartArt or vector-drawn equations may be missing. The compiler
should treat this output as lossy and flag anything uncertain.

Usage:
    python scripts/extract_pptx.py <input_dir> [output_root]
      input_dir   : folder containing *.pptx
      output_root : where to write <stem>/auto/<stem>.md (+ images/);
                    default = <input_dir>/../mineru  so it sits beside mineru output

Note: python-pptx's math->latex path contains a print() that crashes on a GBK
console. This script reconfigures stdout to UTF-8 to avoid that; if you invoke it
yourself, also set $env:PYTHONIOENCODING = 'utf-8'.
"""
import os, sys, glob
try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")
except AttributeError:
    pass
from pptx import Presentation


def shape_text(shape):
    chunks = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = "".join(run.text for run in para.runs).strip()
            if t:
                indent = "  " * para.level
                chunks.append(f"{indent}- {t}" if para.level else t)
    if shape.has_table:
        for row in shape.table.rows:
            chunks.append("| " + " | ".join(c.text.strip() for c in row.cells) + " |")
    return chunks


def extract(pptx_path, out_root):
    stem = os.path.splitext(os.path.basename(pptx_path))[0]
    out_dir = os.path.join(out_root, stem, "auto")
    img_dir = os.path.join(out_dir, "images")
    os.makedirs(img_dir, exist_ok=True)
    prs = Presentation(pptx_path)

    lines = [
        f"# {stem}",
        "",
        "> Source: python-pptx text extraction (**NOT OCR**); layout/positioning lost, "
        "vector math / SmartArt may be missing. Images dumped to images/ but their "
        "relation to the text is not preserved. Treat as lossy; verify against the slide.",
        "",
    ]
    img_n = 0
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n## Slide {i}\n")
        body = []
        for shape in slide.shapes:
            body.extend(shape_text(shape))
            if getattr(shape, "shape_type", None) == 13 or getattr(shape, "image", None):
                try:
                    img = shape.image
                    img_n += 1
                    fn = f"slide{i:03d}_{img_n:03d}.{img.ext or 'png'}"
                    with open(os.path.join(img_dir, fn), "wb") as f:
                        f.write(img.blob)
                    body.append(f"![](images/{fn})")
                except Exception:
                    pass
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                body.append(f"\n_[notes]_ {note}")
        lines.extend(body if body else ["_(no extractable text on this slide)_"])

    md_path = os.path.join(out_dir, f"{stem}.md")
    with open(md_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))
    print(f"  {stem}: {len(prs.slides)} slides, {img_n} images -> {md_path}")


def main():
    if len(sys.argv) < 2:
        print("usage: python scripts/extract_pptx.py <input_dir> [output_root]")
        return 2
    in_dir = os.path.abspath(sys.argv[1])
    out_root = os.path.abspath(sys.argv[2]) if len(sys.argv) > 2 \
        else os.path.join(os.path.dirname(in_dir), "mineru")
    pptx_files = sorted(glob.glob(os.path.join(in_dir, "*.pptx")))
    if not pptx_files:
        print(f"no .pptx in {in_dir}")
        return 0
    print(f"extracting {len(pptx_files)} pptx -> {out_root}")
    for p in pptx_files:
        extract(p, out_root)
    print("done")
    return 0


if __name__ == "__main__":
    sys.exit(main())
